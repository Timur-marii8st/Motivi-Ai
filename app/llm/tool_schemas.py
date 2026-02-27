"""
Tool/function definitions for OpenAI/OpenRouter function calling.
"""

def _to_openai_tool(schema: dict) -> dict:
    """Wraps a function schema in OpenAI's tool format."""
    return {
        "type": "function",
        "function": schema
    }

TOOL_SCHEDULE_REMINDER = {
    "name": "schedule_reminder",
    "description": "Schedule a one-off motivational reminder message for the user at a specific datetime. The time should be in user's local timezone.",
    "parameters": {
        "type": "object",
        "properties": {
            "message_text": {"type": "string", "description": "Text of the reminder message to send"},
            "reminder_datetime_iso": {"type": "string", "description": "Exact datetime for the reminder in ISO format (YYYY-MM-DDTHH:MM:SS). Use the user's local time. Example: 2025-11-26T15:30:00"},
        },
        "required": ["message_text", "reminder_datetime_iso"],
    },
}

TOOL_CANCEL_REMINDER = {
    "name": "cancel_reminder",
    "description": "Cancel a scheduled reminder created with schedule_reminder tool.",
    "parameters": {
        "type": "object",
        "properties": {
            "job_id": {"type": "string", "description": "The job_id returned when the reminder was scheduled"},
        },
        "required": ["job_id"],
    },
}

TOOL_LIST_REMINDERS = {
    "name": "list_reminders",
    "description": "List all active scheduled reminders for the user.",
    "parameters": {
        "type": "object",
        "properties": {},
        "required": [],
    },
}

TOOL_CREATE_PLAN = {
    "name": "create_plan",
    "description": "Create a plan for the user (daily, weekly, or monthly) and send it as a message. The plan will be stored in memory for the specified duration.",
    "parameters": {
        "type": "object",
        "properties": {
            "plan_level": {
                "type": "string",
                "enum": ["daily", "weekly", "monthly"],
                "description": "The time scope of the plan: 'daily' (stored for 1 day), 'weekly' (stored for 7 days), or 'monthly' (stored for 30 days)"
            },
            "plan_content": {"type": "string", "description": "The plan content to send to the user and store in memory"},
        },
        "required": ["plan_level", "plan_content"],
    },
}

TOOL_CHECK_PLAN = {
    "name": "check_plan",
    "description": "Check which active plans are currently stored for the user.",
    "parameters": {
        "type": "object",
        "properties": {},
        "required": [],
    },
}

TOOL_EDIT_PLAN = {
    "name": "edit_plan",
    "description": "Edit an existing plan that was created with create_plan. You can update the plan content and/or extend its expiration time.",
    "parameters": {
        "type": "object",
        "properties": {
            "plan_id": {"type": "integer", "description": "The ID of the plan to edit"},
            "plan_content": {"type": "string", "description": "The new plan content to send to the user and store"},
            "extend_expiry": {
                "type": "boolean",
                "description": "If true, extends the expiry date based on the original plan level from now. Default is false."
            },
        },
        "required": ["plan_id", "plan_content"],
    },
}

TOOL_CREATE_CALENDAR_EVENT = {
    "name": "create_calendar_event",
    "description": "Create an event in the user's Google Calendar.",
    "parameters": {
        "type": "object",
        "properties": {
            "summary": {"type": "string", "description": "Event title"},
            "start_datetime": {"type": "string", "description": "Start datetime ISO format"},
            "end_datetime": {"type": "string", "description": "End datetime ISO format"},
            "description": {"type": "string", "description": "Event description"},
        },
        "required": ["summary", "start_datetime", "end_datetime"],
    },
}

TOOL_CHECK_AVAILABILITY = {
    "name": "check_calendar_availability",
    "description": "Check if the user is available during a time window.",
    "parameters": {
        "type": "object",
        "properties": {
            "start_datetime": {"type": "string", "description": "Start datetime ISO format"},
            "end_datetime": {"type": "string", "description": "End datetime ISO format"},
        },
        "required": ["start_datetime", "end_datetime"],
    },
}


TOOL_EXECUTE_CODE = {
    "name": "execute_code",
    "description": (
        "Execute a code snippet in an isolated sandbox and return the output. "
        "Use this when the user asks you to run code, compute something, generate charts, "
        "create documents (Word/Excel/PowerPoint), or demonstrates a programming concept.\n\n"
        "Supported languages: python, javascript, bash.\n\n"
        "PYTHON FILE OUTPUT — saves files the user receives automatically:\n"
        "  Save any file to /output/ inside the sandbox and it will be sent to the user via Telegram.\n"
        "  Always use /output/ paths — do NOT use plt.show() or open files elsewhere.\n\n"
        "  Pre-installed Python libraries (no pip needed):\n"
        "    matplotlib, numpy, pandas, scipy, seaborn  — charts & data analysis\n"
        "    python-docx   — create Word (.docx) files\n"
        "    openpyxl      — create Excel (.xlsx) spreadsheets\n"
        "    python-pptx   — create PowerPoint (.pptx) presentations\n"
        "    Pillow        — image creation and manipulation\n\n"
        "  Examples:\n"
        "    # Matplotlib chart → PNG sent as photo\n"
        "    import matplotlib.pyplot as plt\n"
        "    plt.plot([1, 2, 3]); plt.savefig('/output/chart.png')\n\n"
        "    # Word document sent as .docx file\n"
        "    from docx import Document\n"
        "    doc = Document(); doc.add_paragraph('Hello'); doc.save('/output/report.docx')\n\n"
        "    # Excel spreadsheet sent as .xlsx file\n"
        "    from openpyxl import Workbook\n"
        "    wb = Workbook(); ws = wb.active; ws['A1'] = 'Data'; wb.save('/output/data.xlsx')\n\n"
        "    # PowerPoint presentation sent as .pptx file\n"
        "    from pptx import Presentation\n"
        "    prs = Presentation(); slide = prs.slides.add_slide(prs.slide_layouts[0])\n"
        "    slide.shapes.title.text = 'My Slide'; prs.save('/output/slides.pptx')\n\n"
        "Sandbox constraints: no network access, 30-second timeout, 256 MB memory, read-only root FS. "
        "Do NOT use for malicious, harmful, or privacy-violating code."
    ),
    "parameters": {
        "type": "object",
        "properties": {
            "language": {
                "type": "string",
                "enum": ["python", "javascript", "bash"],
                "description": "Programming language of the code snippet.",
            },
            "code": {
                "type": "string",
                "description": (
                    "The source code to execute. "
                    "For Python file output, save files to /output/ (e.g. plt.savefig('/output/chart.png')). "
                    "Never call plt.show() — use savefig instead."
                ),
            },
        },
        "required": ["language", "code"],
    },
}

RAW_TOOLS = [
    TOOL_SCHEDULE_REMINDER,
    TOOL_CANCEL_REMINDER,
    TOOL_LIST_REMINDERS,
    TOOL_CREATE_PLAN,
    TOOL_CHECK_PLAN,
    TOOL_EDIT_PLAN,
    TOOL_CREATE_CALENDAR_EVENT,
    TOOL_CHECK_AVAILABILITY,
    TOOL_EXECUTE_CODE,
]

ALL_TOOLS = [_to_openai_tool(t) for t in RAW_TOOLS]